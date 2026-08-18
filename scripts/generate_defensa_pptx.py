#!/usr/bin/env python3
"""Genera presentacion_defensa.pptx para la defensa del TFM."""

from __future__ import annotations

from pathlib import Path

PROJECT_ROOT = Path(__file__).resolve().parents[1]
FIGURES = PROJECT_ROOT / "docs" / "figures"
OUTPUT = PROJECT_ROOT / "docs" / "presentacion_defensa.pptx"

AUTOR = "Ángela Ramírez Barajas"
DIRECTOR = "Juanjo Salvador"
UNIVERSIDAD = "Universidad Católica de Murcia (UCAM)"
MASTER = "2ª Ed. Máster en IA Aplicada a la Ciberseguridad"


def _set_title(slide, title: str) -> None:
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.32), Inches(9), Inches(0.75))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True


def add_title_slide(prs, title: str, subtitle: str = "") -> None:
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.15), Inches(8.8), Inches(1.8))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    box.text_frame.word_wrap = True

    if subtitle:
        box2 = slide.shapes.add_textbox(Inches(0.6), Inches(3.2), Inches(8.8), Inches(3.2))
        tf = box2.text_frame
        tf.word_wrap = True
        for i, line in enumerate(subtitle.split("\n")):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = line
            para.font.size = Pt(18)


def add_bullet_slide(prs, title: str, bullets: list[str], image: Path | None = None) -> None:
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_title(slide, title)

    body_width = Inches(4.6) if image and image.exists() else Inches(8.8)
    body = slide.shapes.add_textbox(Inches(0.55), Inches(1.15), body_width, Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(15)
        p.space_after = Pt(8)

    if image and image.exists():
        slide.shapes.add_picture(str(image), Inches(5.35), Inches(1.25), width=Inches(4.2))


def main() -> None:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise SystemExit("pip install python-pptx") from exc

    prs = Presentation()
    prs.slide_width = 9144000
    prs.slide_height = 6858000

    add_title_slide(
        prs,
        "Detección de secretos en código y pipelines mediante reglas y LLM",
        f"{AUTOR}\n{MASTER}\n{UNIVERSIDAD}\nDirector: {DIRECTOR}\nJulio 2026",
    )

    slides: list[tuple[str, list[str], Path | None]] = [
        (
            "El problema (y el enunciado)",
            [
                "Las reglas (regex + entropía) son rápidas e integrables en Git",
                "El coste es el ruido: un mock o un README dispara la misma alerta que un secreto real",
                "Enunciado: reducir FP, integrar en pre-commit/CI, medir, rotar y ocultar",
            ],
            None,
        ),
        (
            "Cómo lo hago: dos capas",
            [
                "Gitleaks detecta candidatos (única capa que encuentra secretos)",
                "Cruce con CredData: TP (6.845) y FP (1.128)",
                "El LLM local (Llama 3.1 8B) solo ve FP: decide si descartar la alerta",
                "Precisión híbrida = TP / (TP + FP que el filtro no descartó)",
            ],
            FIGURES / "fig01_pipeline_hibrido.png",
        ),
        (
            "Por qué solo evalúo falsos positivos",
            [
                "El LLM no es un detector: no busca secretos en el código",
                "Si viera también TP, un error suyo bajaría el recall (silenciaría un secreto real)",
                "El enunciado pide reducir ruido, no sustituir las reglas",
                "Por eso el recall del híbrido es exactamente el de Gitleaks: 45,86 %",
            ],
            None,
        ),
        (
            "¿Mejora a otras propuestas?",
            [
                "Frente a reglas solas: sí. Precisión 85,85 % → 99,97 % proyectada; mismo recall",
                "Frente a Rahman et al. / Biringa (LLM o ML con entrenamiento): no reclamo mejor F1",
                "Ellos maximizan F1 global; yo filtro FP de Gitleaks sin fine-tuning",
                "Aporte: operativo, local (Ollama), y el prompt contextual es lo que marca la diferencia",
            ],
            None,
        ),
        (
            "Baseline: solo Gitleaks",
            [
                "CredData: ~1,02 GB en 56,8 s | 8.210 hallazgos",
                "TP 6.845 | FP 1.128 | FN 8.177",
                "Precisión 85,85 % | Recall 45,86 % | F1 59,79 %",
                "Útil para detectar; inviable de triar a mano en un hook diario",
            ],
            FIGURES / "fig05_precision_recall.png",
        ),
        (
            "El prompt es el resultado",
            [
                "Mismo modelo, temperatura 0, mismos 200 FP",
                "v1 genérico («¿parece un secreto?»): 34 % de acierto",
                "v2 contextual (rutas test/mock/docs + nombres MOCK): 99 %",
                "130 correcciones, 0 regresiones, 1 caso residual (PEM en docs inline)",
            ],
            FIGURES / "fig02_precision_comparativa.png",
        ),
        (
            "Cómo se opera",
            [
                "Pre-commit y GitHub Actions: solo Gitleaks (bloqueo en < 1 s)",
                "LLM en lote o runner con Ollama (~5–15 s/candidato)",
                "Secretos de pipeline: referencia a secrets del proveedor, nunca en el YAML",
                "Hallazgo real en main: rotar en < 24 h; allowlist solo con justificación",
            ],
            None,
        ),
        (
            "Limitaciones",
            [
                "Muestra: 200 de 1.128 FP (proyección 99,97 % = cota alta)",
                "El recall de las reglas no mejora: el LLM no ve los FN",
                "Un modelo local; sin verificar si la credencial está viva",
                "v2 puede estar alineado a patrones típicos de CredData",
            ],
            None,
        ),
        (
            "Conclusiones",
            [
                "Híbrido viable sin fine-tuning ni cloud: reglas detectan, LLM tria FP",
                "H1 confirmada: 99 % de FP descartados en la muestra (umbral 30 %)",
                "La diferencia no es «añadir un LLM», es el prompt de contexto",
                "Pre-commit + CI con reglas; políticas de rotación y ocultación",
            ],
            None,
        ),
        (
            "¿Preguntas?",
            [
                f"{AUTOR} — Director: {DIRECTOR}",
                f"{MASTER}",
                UNIVERSIDAD,
            ],
            None,
        ),
    ]

    for title, bullets, image in slides:
        add_bullet_slide(prs, title, bullets, image)

    prs.save(OUTPUT)
    print(f"Presentacion guardada en: {OUTPUT}")


if __name__ == "__main__":
    main()
